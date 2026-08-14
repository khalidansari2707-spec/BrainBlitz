import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_ppt():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800
    CARD_BORDER = RGBColor(51, 65, 85)   # Slate 700
    TEXT_WHITE = RGBColor(248, 250, 252) # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ACCENT_CYAN = RGBColor(56, 189, 248) # Sky 400
    ACCENT_INDIGO = RGBColor(129, 140, 248) # Indigo 400
    ACCENT_PINK = RGBColor(244, 63, 94)  # Rose 500
    ACCENT_GREEN = RGBColor(52, 211, 153) # Emerald 400

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="BRAINBLITZ PRESENTATION"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        p_cat.font.name = "Arial"

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)

    # Top accent bar
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.2), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_CYAN
    accent_bar.line.fill.background()

    # Main Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "BrainBlitz ⚡"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Next-Gen Interactive Educational Quiz & Diagram Game"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_CYAN
    p2.font.name = "Arial"
    p2.space_before = Pt(10)

    # Subtitle Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.733), Inches(3.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER
    card1.line.width = Pt(1)

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = Inches(0.4)
    tf_c1.margin_top = Inches(0.3)
    tf_c1.margin_right = Inches(0.4)

    items_s1 = [
        ("🚀 Live Application URL:", "https://brainblitz-23d97.web.app/", ACCENT_GREEN),
        ("🤖 Powered by AI Toolchain:", "NVIDIA Nemotron 3.5 (via OpenRouter) & Cline Agent in VS Code", ACCENT_INDIGO),
        ("📚 Multi-Disciplinary Learning:", "Biology, Chemistry, Physics, OS, DBMS, Computer Networks & Software Eng.", TEXT_WHITE),
        ("💡 Developed By:", "Khalid Ansari", ACCENT_CYAN)
    ]

    for label, val, color in items_s1:
        p_lbl = tf_c1.add_paragraph()
        p_lbl.text = label + " "
        p_lbl.font.bold = True
        p_lbl.font.size = Pt(15)
        p_lbl.font.color.rgb = color
        p_lbl.font.name = "Arial"

        # Add value part
        p_val = p_lbl.add_run()
        p_val.text = val
        p_val.font.bold = False
        p_val.font.size = Pt(15)
        p_val.font.color.rgb = TEXT_WHITE
        p_lbl.space_after = Pt(12)

    # ==========================================
    # SLIDE 2: Project Overview & Core Features
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "Project Summary & Core Game Modes", "PROJECT OVERVIEW")

    # 4 Grid Feature Cards
    card_width = Inches(5.6)
    card_height = Inches(2.3)

    features = [
        ("🎓 7 STEM & Computer Subjects", "Comprehensive interactive quizzes spanning Biology, Chemistry, Physics, OS, DBMS, Computer Networks, and Software Engineering.", ACCENT_CYAN, Inches(0.8), Inches(1.8)),
        ("⚡ 5 Engaging Game Modes", "Includes Label Rush (Drag & Drop), Sequence Snap (Workflow ordering), Spot the Fault (Error finding), Match Links, & Speed Blitz (30s MCQ).", ACCENT_INDIGO, Inches(6.9), Inches(1.8)),
        ("🎨 Modern Tech & UI/UX", "Built with React 19, Vite, Tailwind CSS, Lucide Icons, and Canvas Confetti for gamified feedback and dynamic micro-animations.", ACCENT_GREEN, Inches(0.8), Inches(4.4)),
        ("🚀 Production Cloud Deployment", "Fully optimized production build deployed to Firebase Hosting with global low-latency CDN distribution.", ACCENT_PINK, Inches(6.9), Inches(4.4)),
    ]

    for title, desc, accent_color, pos_x, pos_y in features:
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        tf_f = card.text_frame
        tf_f.word_wrap = True
        tf_f.margin_left = Inches(0.3)
        tf_f.margin_top = Inches(0.25)
        tf_f.margin_right = Inches(0.3)

        pf_t = tf_f.paragraphs[0]
        pf_t.text = title
        pf_t.font.bold = True
        pf_t.font.size = Pt(17)
        pf_t.font.color.rgb = accent_color
        pf_t.font.name = "Arial"

        pf_d = tf_f.add_paragraph()
        pf_d.text = desc
        pf_d.font.size = Pt(13)
        pf_d.font.color.rgb = TEXT_MUTED
        pf_d.font.name = "Arial"
        pf_d.space_before = Pt(8)

    # ==========================================
    # SLIDE 3: AI Toolchain & Tech Stack
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "AI Development Stack & Tech Infrastructure", "AI & TECH STACK")

    # 3 Column Cards
    col_width = Inches(3.64)
    col_height = Inches(4.9)

    cols = [
        ("🤖 Nemotron 3.5 AI Model", [
            ("Model:", "NVIDIA Nemotron 3.5"),
            ("Role:", "High-Reasoning LLM for code logic, component architecture, and complex game state design."),
            ("Capabilities:", "Exceptional structured JSON output generation, SVG diagram mapping, and fast context handling.")
        ], ACCENT_CYAN, Inches(0.8)),
        ("⚡ OpenRouter & Cline Agent", [
            ("Environment:", "VS Code + Cline Agent Extension"),
            ("Gateway:", "OpenRouter API Platform"),
            ("Workflow:", "Autonomous multi-file editing, terminal command execution, and rapid iterative debugging."),
            ("Efficiency:", "Accelerated development lifecycle from hours to minutes.")
        ], ACCENT_INDIGO, Inches(4.84)),
        ("🌐 Web & Cloud Stack", [
            ("Frontend Framework:", "React 19 + Vite"),
            ("Styling & Design:", "Tailwind CSS + Glassmorphism"),
            ("Animations:", "Canvas Confetti"),
            ("Hosting & CDN:", "Firebase Hosting"),
            ("Version Control:", "Git & GitHub Repository")
        ], ACCENT_GREEN, Inches(8.88))
    ]

    for title, bullet_list, accent_color, pos_x in cols:
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, Inches(1.8), col_width, col_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_top = Inches(0.25)
        tf_c.margin_right = Inches(0.25)

        p_ct = tf_c.paragraphs[0]
        p_ct.text = title
        p_ct.font.bold = True
        p_ct.font.size = Pt(17)
        p_ct.font.color.rgb = accent_color
        p_ct.font.name = "Arial"

        for label, text in bullet_list:
            p_b = tf_c.add_paragraph()
            p_b.text = "• " + label + " "
            p_b.font.bold = True
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = TEXT_WHITE
            p_b.font.name = "Arial"
            p_b.space_before = Pt(8)

            r_t = p_b.add_run()
            r_t.text = text
            r_t.font.bold = False
            r_t.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 4: How Nemotron 3.5 Was Accessed & Integrated
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "How Nemotron 3.5 Was Accessed via OpenRouter & Cline", "AI INTEGRATION ARCHITECTURE")

    # Step Flow - 4 horizontal cards
    step_width = Inches(11.733)
    step_height = Inches(1.05)

    steps = [
        ("Step 1: OpenRouter API Key Setup", "Obtained API Key from OpenRouter platform enabling unified access to NVIDIA Nemotron 3.5 LLM endpoints.", ACCENT_CYAN, Inches(1.8)),
        ("Step 2: Cline Extension Configuration in VS Code", "Configured VS Code's Cline Agent with OpenRouter provider settings, setting Model ID to 'nvidia/nemotron-3.5'.", ACCENT_INDIGO, Inches(3.05)),
        ("Step 3: Prompt Engineering & Prompt Routing", "Formulated high-context system prompts for interactive UI components, state management, and diagram data schema.", ACCENT_GREEN, Inches(4.3)),
        ("Step 4: Autonomous Execution & Deployment", "Cline + Nemotron 3.5 generated React components, debugged logic, built dist bundle, and executed Firebase deployment.", ACCENT_PINK, Inches(5.55)),
    ]

    for title, text, accent_color, pos_y in steps:
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), pos_y, step_width, step_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        tf_s = card.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = Inches(0.3)
        tf_s.margin_top = Inches(0.18)

        p_st = tf_s.paragraphs[0]
        p_st.text = title
        p_st.font.bold = True
        p_st.font.size = Pt(15)
        p_st.font.color.rgb = accent_color
        p_st.font.name = "Arial"

        p_sd = tf_s.add_paragraph()
        p_sd.text = text
        p_sd.font.size = Pt(12)
        p_sd.font.color.rgb = TEXT_MUTED
        p_sd.font.name = "Arial"
        p_sd.space_before = Pt(3)

    # ==========================================
    # SLIDE 5: Live Deployment & Future Roadmap
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "Live Web Deployment & Future AI Roadmap", "DEPLOYMENT & FUTURE")

    # Left Column: Deployment Info
    card_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = CARD_BORDER
    card_l.line.width = Pt(1)

    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    tf_l.margin_right = Inches(0.3)

    p_lt = tf_l.paragraphs[0]
    p_lt.text = "🌐 Live Production Status"
    p_lt.font.bold = True
    p_lt.font.size = Pt(18)
    p_lt.font.color.rgb = ACCENT_GREEN
    p_lt.font.name = "Arial"

    dep_details = [
        ("Hosting Platform:", "Firebase Hosting (Google Cloud)"),
        ("Live URL:", "https://brainblitz-23d97.web.app/"),
        ("Build System:", "Vite 8.2 + Rollup production bundle"),
        ("Global Latency:", "< 50ms worldwide via Firebase CDN"),
        ("Status:", "Live & Fully Operational ✅")
    ]

    for label, val in dep_details:
        p_d = tf_l.add_paragraph()
        p_d.text = "• " + label + " "
        p_d.font.bold = True
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(10)

        r_v = p_d.add_run()
        r_v.text = val
        r_v.font.bold = False
        r_v.font.color.rgb = ACCENT_CYAN if "http" in val else TEXT_MUTED

    # Right Column: Future Roadmap
    card_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.633), Inches(4.9))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = CARD_BORDER
    card_r.line.width = Pt(1)

    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    tf_r.margin_right = Inches(0.3)

    p_rt = tf_r.paragraphs[0]
    p_rt.text = "🔮 Future Enhancements Roadmap"
    p_rt.font.bold = True
    p_rt.font.size = Pt(18)
    p_rt.font.color.rgb = ACCENT_INDIGO
    p_rt.font.name = "Arial"

    roadmap = [
        ("🏆 Global Leaderboard", "Integration with Firebase Firestore & Authentication for real-time multiplayer scoring."),
        ("🧠 Dynamic AI Quiz Generator", "Real-time question generation powered by Nemotron 3.5 API based on user topic input."),
        ("📱 PWA & Mobile App Support", "Offline quiz capability and installable Progressive Web Application."),
        ("📊 Analytics Dashboard", "Detailed performance tracking & mastery reports for students and educators.")
    ]

    for rm_title, rm_desc in roadmap:
        p_rm = tf_r.add_paragraph()
        p_rm.text = rm_title
        p_rm.font.bold = True
        p_rm.font.size = Pt(13)
        p_rm.font.color.rgb = TEXT_WHITE
        p_rm.space_before = Pt(8)

        p_rmd = tf_r.add_paragraph()
        p_rmd.text = rm_desc
        p_rmd.font.size = Pt(11)
        p_rmd.font.color.rgb = TEXT_MUTED
        p_rmd.space_before = Pt(2)

    # Save presentation
    output_path = "c:\\Users\\Salman Ansari\\Desktop\\First Game\\first-game\\BrainBlitz_Presentation.pptx"
    prs.save(output_path)
    print(f"PPT successfully created at: {output_path}")

if __name__ == "__main__":
    build_ppt()
